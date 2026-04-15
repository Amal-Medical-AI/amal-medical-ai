import streamlit as st
from groq import Groq
import PyPDF2
from docx import Document
from pptx import Presentation
import re

# 1. رجعنا الاسم والشعار الخاص فيكي
st.set_page_config(page_title="Amal's Medical Brain", layout="wide")
st.markdown("""
    <style>
    .main { background-color: #fff0f5; }
    .stButton>button { width: 100%; border-radius: 20px; background-color: #ff69b4; color: white; font-weight: bold; margin-bottom: 10px; height: 3em; border: none; }
    .stButton>button:hover { background-color: #d81b60; transform: scale(1.02); }
    .master-box { background-color: white; padding: 20px; border-radius: 15px; border: 2px solid #ff69b4; color: #333; }
    h1, h3 { color: #ff69b4; text-align: center; }
    </style>
    """, unsafe_allow_html=True)

client = Groq(api_key=st.secrets["GROQ_API_KEY"])

# 2. نظام الذاكرة التراكمية (عشان ما يضيع شي)
if "db" not in st.session_state:
    st.session_state.db = {"General": {"summary": "", "questions": ""}}

# 3. دالة قراءة كل أنواع الملفات (PDF, Word, PPTX)
def get_text(file):
    ext = file.name.split('.')[-1].lower()
    t = ""
    try:
        if ext == 'pdf':
            r = PyPDF2.PdfReader(file)
            for p in r.pages[:8]: t += p.extract_text() + " "
        elif ext == 'docx':
            d = Document(file)
            for p in d.paragraphs: t += p.text + " "
        elif ext in ['pptx', 'ppt']:
            prs = Presentation(file)
            for s in prs.slides:
                for shp in s.shapes:
                    if hasattr(shp, "text"): t += shp.text + " "
    except: return "Error reading file."
    return t

# --- القائمة الجانبية (إدارة الكورسات باسمك) ---
st.sidebar.title("🌸 Amal's Courses")
c_name = st.sidebar.text_input("New Course (e.g. Pathology):")
if st.sidebar.button("Add Course"):
    if c_name and c_name not in st.session_state.db:
        st.session_state.db[c_name] = {"summary": "", "questions": ""}

active_c = st.sidebar.selectbox("Current Course:", list(st.session_state.db.keys()))

# --- الصفحة الرئيسية ---
st.title("🌸 Amal's Medical Brain 🌸")
st.subheader(f"📍 Study Lab: {active_c}")

uploaded_file = st.file_uploader("Upload (GoodNotes PDF, Word, PowerPoint)", type=["pdf", "docx", "pptx"])

if uploaded_file:
    with st.spinner("Processing High-Speed Analysis..."):
        clean_text = re.sub(r'\s+', ' ', get_text(uploaded_file)).strip()[:6000]

    def ask_ai(prompt):
        resp = client.chat.completions.create(messages=[{"role":"user","content":prompt}], model="llama-3.3-70b-versatile")
        return resp.choices[0].message.content

    # --- الأزرار التراكمية (ورا بعض) ---
    st.write("---")
    
    if st.button("📝 Summarize & Add to Course Master File"):
        with st.spinner('Updating internal memory...'):
            p = f"Previous Master Summary: {st.session_state.db[active_c]['summary']}\nNew Info: {clean_text}\nUpdate the unified medical summary."
            st.session_state.db[active_c]['summary'] = ask_ai(p)
            st.success("Summary added to Master Record!")

    if st.button("❓ Create & Save USMLE Questions"):
        with st.spinner('Building your Question Bank...'):
            p = f"Current Bank: {st.session_state.db[active_c]['questions']}\nNew Text: {clean_text}\nAdd 3 USMLE questions to the bank."
            st.session_state.db[active_c]['questions'] += "\n\n" + ask_ai(p)
            st.success("Questions added to Bank!")

    # --- عرض مجهودك التراكمي ---
    st.divider()
    with st.expander(f"📚 View {active_c} Master Summary"):
        st.markdown(f'<div class="master-box">{st.session_state.db[active_c]["summary"]}</div>', unsafe_allow_html=True)
    
    with st.expander(f"❓ View {active_c} Question Bank"):
        st.markdown(f'<div class="master-box">{st.session_state.db[active_c]["questions"]}</div>', unsafe_allow_html=True)

    # زر التحميل لـ OneDrive
    full_data = f"AMAL'S MASTER STUDY GUIDE - {active_c}\n\nSUMMARY:\n{st.session_state.db[active_c]['summary']}\n\nQUESTIONS:\n{st.session_state.db[active_c]['questions']}"
    st.download_button(label="📥 Download Everything for OneDrive", data=full_data, file_name=f"Amal_{active_c}_Master.txt")

# الترجمة (جانبية)
st.sidebar.divider()
if st.sidebar.button("🌐 Translate Summary to Arabic"):
    st.write(ask_ai("Translate this to Arabic: " + st.session_state.db[active_c]['summary'][:2000]))

st.divider()
user_q = st.text_input("💬 Ask Amal's AI anything:")
if user_q: st.write(ask_ai(user_q))
